#!/usr/bin/env python3

import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Book cover inspired color scheme
COLORS = {
    'primary': RGBColor(133, 153, 164),    # Steel blue-gray
    'secondary': RGBColor(255, 255, 255),   # White
    'accent': RGBColor(152, 169, 179),      # Lighter steel blue
    'background': RGBColor(255, 255, 255),  # White background
    'text': RGBColor(36, 41, 46),          # Dark text
    'highlight': RGBColor(144, 162, 172)    # Medium steel blue
}

# Font settings
FONTS = {
    'title': 'Helvetica Neue',
    'heading': 'Helvetica Neue',
    'body': 'Helvetica Neue'
}

def check_dependencies():
    """Check if required dependencies are installed."""
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx...")
        import subprocess
        import sys
        subprocess.run([sys.executable, "-m", "pip", "install", "-r", "requirements.txt"], check=True)

def create_output_dirs():
    """Create necessary output directories."""
    Path("output").mkdir(parents=True, exist_ok=True)

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def read_markdown_content(file_path):
    """Read and parse markdown content file."""
    with open(file_path, 'r') as f:
        content = f.read()
    
    # Split into slides based on '---' separator
    slides = content.split('---')
    
    # Parse title and subtitle from first slide
    title_slide = slides[0].strip().split('\n')
    title = title_slide[0].lstrip('#').strip()
    subtitle = '\n'.join(title_slide[1:]).strip()
    
    # Parse content slides
    content_slides = []
    for slide in slides[1:]:
        lines = slide.strip().split('\n')
        if not lines:
            continue
            
        # First line is the heading
        heading = lines[0].lstrip('#').strip()
        
        # Rest are bullet points
        bullets = []
        for line in lines[1:]:
            if line.strip():
                # Remove markdown bullet points and numbering
                bullet = re.sub(r'^[\d\.\-\*]+', '', line).strip()
                bullets.append(bullet)
        
        content_slides.append((heading, bullets))
    
    return title, subtitle, content_slides

def create_title_slide(prs, title, subtitle):
    """Create the title slide with book-inspired design."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Set background
    set_slide_background(slide, COLORS['background'])
    
    # Add subtle header bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(16)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Add book cover image
    try:
        image_path = Path("slides/assets/images/book-cover.jpg")
        if image_path.exists():
            left = Inches(4)
            top = Inches(1)
            width = Inches(8)
            height = Inches(4)
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            
            # Adjust title position to be below the image
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.name = FONTS['title']
            title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
            title_shape.text_frame.paragraphs[0].font.size = Pt(48)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.top = Inches(5.5)
            
            # Adjust subtitle position
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.name = FONTS['body']
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.top = Inches(7)
        else:
            print("Warning: Book cover image not found at", image_path)
            # Use default title slide layout if image is not found
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.name = FONTS['title']
            title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
            title_shape.text_frame.paragraphs[0].font.size = Pt(48)
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.name = FONTS['body']
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    except Exception as e:
        print(f"Error adding book cover image: {e}")
        # Fall back to default title slide if there's an error
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.name = FONTS['title']
        title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
        title_shape.text_frame.paragraphs[0].font.size = Pt(48)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.name = FONTS['body']
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

def create_section_title_slide(prs, title):
    """Create a section title slide with book-inspired design."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background
    set_slide_background(slide, COLORS['background'])
    
    # Add subtle header
    left = Inches(0)
    top = Inches(0)
    width = Inches(16)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Add title
    left = Inches(1)
    top = Inches(2)
    width = Inches(14)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = title
    p.font.name = FONTS['heading']
    p.font.size = Pt(42)
    p.font.color.rgb = COLORS['primary']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtle decorative line
    left = Inches(4)
    top = Inches(4)
    width = Inches(8)
    height = Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['accent']

def create_content_slide(prs, title, content):
    """Create a content slide with book-inspired design."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set background
    set_slide_background(slide, COLORS['background'])
    
    # Add subtle header
    left = Inches(0)
    top = Inches(0)
    width = Inches(16)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.7)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = title
    p.font.name = FONTS['heading']
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['primary']
    p.font.bold = True
    
    # Add content
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FONTS['body']
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text']
        p.level = 0
        p.space_after = Pt(12)
    
    # Add subtle footer
    left = Inches(0)
    top = Inches(8.7)
    width = Inches(16)
    height = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']

def generate_slides():
    """Generate PowerPoint presentation with book-inspired design."""
    print("Generating slides...")
    
    # Read content from markdown file
    content_file = Path("slides/content/01-introduction.md")
    if not content_file.exists():
        print(f"Error: Content file not found at {content_file}")
        return
    
    title, subtitle, content_slides = read_markdown_content(content_file)
    
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create slides
    create_title_slide(prs, title, subtitle)
    
    # Create content slides
    for heading, bullets in content_slides:
        create_section_title_slide(prs, heading)
        create_content_slide(prs, heading, bullets)
    
    # Save the presentation
    output_path = Path("output/microprediction-update.pptx")
    prs.save(str(output_path))
    print(f"\nPresentation saved to: {output_path}")

def main():
    """Main function to orchestrate slide generation."""
    print("Starting slide generation...")
    
    # Check and install dependencies
    check_dependencies()
    
    # Create output directories
    create_output_dirs()
    
    # Generate slides
    generate_slides()
    
    print("\nSlide generation complete!")

if __name__ == "__main__":
    main() 